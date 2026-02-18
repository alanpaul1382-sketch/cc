"""
Script to add chart titles and axis labels to Slides.pptx.
Requires: pip install python-pptx

Usage: python update_charts.py
"""
from pptx import Presentation


def update_charts(input_file="Slides.pptx", output_file="Slides.pptx"):
    prs = Presentation(input_file)

    # Chart updates: slide number -> {title, x_axis, y_axis}
    # None means the label already exists or is not applicable
    chart_updates = {
        7: {
            "title": "EdTech Market Size Comparison",
            "x_axis": None,  # Already has 'Market'
            "y_axis": None,  # Already has 'Market Size (Bilions)'
        },
        10: {
            "title": "AI Cost-Benefit Analysis",
            "x_axis": "Year",
            "y_axis": "SGD (Millions)",
        },
        12: {
            "title": "VR/AR Cost-Benefit Analysis",
            "x_axis": "Year",
            "y_axis": "SGD (Millions)",
        },
        14: {
            "title": "5G & Edge Computing Cost-Benefit Analysis",
            "x_axis": "Year",
            "y_axis": "SGD (Millions)",
        },
        16: {
            "title": "Blockchain Cost-Benefit Analysis",
            "x_axis": "Year",
            "y_axis": "SGD (Millions)",
        },
        # Slide 18: Pie chart — already has title, no axes
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slide_num not in chart_updates:
            continue

        updates = chart_updates[slide_num]

        for shape in slide.shapes:
            if not shape.has_chart:
                continue

            chart = shape.chart

            # Add title if needed
            if updates.get("title") and not chart.has_title:
                chart.has_title = True
                chart.chart_title.has_text_frame = True
                chart.chart_title.text_frame.text = updates["title"]
                print(f"Slide {slide_num}: Added title '{updates['title']}'")

            # Add X axis (category axis) label if needed
            if updates.get("x_axis"):
                try:
                    cat_axis = chart.category_axis
                    if not cat_axis.has_title:
                        cat_axis.has_title = True
                        cat_axis.axis_title.has_text_frame = True
                        cat_axis.axis_title.text_frame.text = updates["x_axis"]
                        print(f"Slide {slide_num}: Added X axis '{updates['x_axis']}'")
                except Exception as e:
                    print(f"Slide {slide_num}: Cannot set X axis — {e}")

            # Add Y axis (value axis) label if needed
            if updates.get("y_axis"):
                try:
                    val_axis = chart.value_axis
                    if not val_axis.has_title:
                        val_axis.has_title = True
                        val_axis.axis_title.has_text_frame = True
                        val_axis.axis_title.text_frame.text = updates["y_axis"]
                        print(f"Slide {slide_num}: Added Y axis '{updates['y_axis']}'")
                except Exception as e:
                    print(f"Slide {slide_num}: Cannot set Y axis — {e}")

    prs.save(output_file)
    print(f"\nSaved to {output_file}")


if __name__ == "__main__":
    update_charts()
