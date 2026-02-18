"""
Generate a PowerPoint presentation deck covering the key points from final.docx.
Topic: Emerging Technologies for Nanyang Polytechnic (NYP)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── colour palette ──────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5C)   # NYP-inspired dark blue
ACCENT     = RGBColor(0x00, 0x7A, 0xCC)   # bright blue accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
MED_GREY   = RGBColor(0x66, 0x66, 0x66)
GREEN      = RGBColor(0x00, 0x80, 0x60)
ORANGE     = RGBColor(0xE8, 0x6C, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── helper functions ────────────────────────────────────────────────────────
def add_bg(slide, colour):
    """Fill the entire slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, colour, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, colour=DARK_GREY, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items,
                     font_size=16, colour=DARK_GREY, spacing=Pt(6)):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


# ── slide builder helpers ───────────────────────────────────────────────────
def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    # accent bar
    add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                title, font_size=40, bold=True, colour=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(2),
                subtitle, font_size=20, colour=RGBColor(0xBB, 0xCC, 0xDD),
                alignment=PP_ALIGN.CENTER)


def make_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ACCENT)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                title, font_size=36, bold=True, colour=WHITE,
                alignment=PP_ALIGN.CENTER)


def make_content_slide(title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # top bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    # title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                title, font_size=28, bold=True, colour=DARK_BLUE)
    # bullets
    add_bullet_slide(slide, Inches(1), Inches(1.3), Inches(11), Inches(5.5),
                     bullets, font_size=18, colour=DARK_GREY)
    return slide


def make_two_column_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                title, font_size=28, bold=True, colour=DARK_BLUE)
    # left column
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
                left_title, font_size=20, bold=True, colour=ACCENT)
    add_bullet_slide(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                     left_items, font_size=16, colour=DARK_GREY)
    # right column
    add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.5),
                right_title, font_size=20, bold=True, colour=ACCENT)
    add_bullet_slide(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5),
                     right_items, font_size=16, colour=DARK_GREY)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE CONTENT
# ═══════════════════════════════════════════════════════════════════════════

# 1 ── Title slide ───────────────────────────────────────────────────────────
make_title_slide(
    "Emerging Technologies for Nanyang Polytechnic",
    "Competency Unit: Emerging Technologies in the Digital Economy\n"
    "Organization: Nanyang Polytechnic (NYP)  •  Sector: Education\n"
    "Team: Clifton (AI) • Collin (VR/AR) • Alex (5G & Edge) • Declan (Blockchain)\n"
    "February 2026"
)

# 2 ── Agenda ────────────────────────────────────────────────────────────────
make_content_slide("Agenda", [
    "1.  Executive Summary",
    "2.  Organization Overview & Value Chain",
    "3.  Current Problems and Challenges",
    "4.  Market Size & Potential Revenue",
    "5.  Technology Solutions",
    "     a) Artificial Intelligence (AI) — Clifton",
    "     b) Immersive Media (VR/AR) — Collin",
    "     c) 5G and Edge Computing — Alex",
    "     d) Blockchain-Based Certificate Verification — Declan",
    "6.  Consolidated Investment & ROI Summary",
    "7.  Three-Phase Implementation Roadmap",
    "8.  Technology Synergy Ecosystem",
    "9.  Conclusion",
])

# 3 ── Executive Summary ────────────────────────────────────────────────────
make_section_slide("Executive Summary")

make_content_slide("Executive Summary", [
    "NYP serves ~15,000 full-time students and thousands of CET learners annually.",
    "Four emerging technologies are proposed to address NYP's pressing challenges:",
    "   • AI — Adaptive learning & predictive analytics",
    "   • VR/AR — Virtual labs & immersive training",
    "   • 5G & Edge Computing — Smart campus backbone",
    "   • Blockchain — Tamper-proof digital credentialing",
    "",
    "Total Year 1 investment: SGD 8.8–17.1 million",
    "Projected annual returns: SGD 10.5–23 million (revenue + cost savings)",
    "Breakeven within 1–2 years for most initiatives.",
    "A phased 3-year implementation roadmap manages risk & maximises synergies.",
])

# 4 ── Organization Overview ────────────────────────────────────────────────
make_section_slide("Organization Overview")

make_content_slide("About Nanyang Polytechnic (NYP)", [
    "Established 1992  •  Located in Ang Mo Kio, Singapore",
    "One of Singapore's five government-funded polytechnics",
    "~15,000 full-time students; thousands of CET learners",
    "6 Schools: IT, Engineering, Business, Design & Media, Health & Social Sciences, Applied Science",
    "Mission: Nurture industry-ready professionals through practice-oriented education",
    "Annual operating budget: ~SGD 200–300 million (MOE-funded)",
])

make_content_slide("NYP's Educational Value Chain", [
    "1.  Student Recruitment & Admissions",
    "2.  Curriculum Design & Development",
    "3.  Teaching & Learning Delivery",
    "4.  Student Support Services",
    "5.  Assessment & Certification",
    "6.  Industry Partnerships & Internships",
    "7.  Continuing Education & Training (CET)",
    "8.  Campus Operations & Infrastructure",
    "",
    "Each value-chain activity is an opportunity for technology-driven improvement.",
])

# 5 ── Current Problems & Challenges ────────────────────────────────────────
make_section_slide("Current Problems & Challenges")

make_content_slide("Key Challenges Facing NYP", [
    "1.  Scalability of Personalized Learning — large class sizes, diverse student needs",
    "2.  Limited Hands-On Training — costly physical labs, limited scheduling, COVID risk",
    "3.  Campus Connectivity Limitations — bandwidth congestion, IoT device growth",
    "4.  Certificate Fraud & Verification Delays — manual process, 3–10 working days",
    "5.  Rising Competition — 5 polytechnics + universities + private institutions",
    "6.  Environmental & Sustainability Concerns — carbon footprint, Singapore Green Plan 2030",
])

# 6 ── Market Size & Revenue ────────────────────────────────────────────────
make_section_slide("Market Size & Potential Revenue")

make_two_column_slide(
    "Market Size & Revenue Opportunity",
    "Education Market",
    [
        "SG govt allocates ~SGD 13B/year to education",
        "Polytechnic sector: ~70,000 students",
        "NYP's share: ~21% of polytechnic market",
        "Global EdTech: USD 340B (2024) → USD 605B (2027)",
        "Asia-Pacific EdTech: USD 115B, ~18% CAGR",
        "SG EdTech: SGD 1.5–2B",
    ],
    "Revenue Streams for NYP",
    [
        "AI-powered CET courses: SGD 5–10M/year",
        "VR/AR B2B content licensing: SGD 2–5M/year",
        "5G/IoT testbed partnerships: SGD 1–2M/year",
        "Blockchain verification services: SGD 0.5–1M/year",
        "Cost savings through automation (AI): SGD 2–4M/year",
        "Smart energy management savings: 10–15%",
    ]
)

# 7 ── AI ───────────────────────────────────────────────────────────────────
make_section_slide("Technology 1: Artificial Intelligence (AI) — Clifton")

make_content_slide("AI — How It Improves NYP", [
    "a) Adaptive Learning Platforms — personalized paths, real-time difficulty adjustment,",
    "     early at-risk student identification (e.g., DreamBox, Knewton)",
    "b) AI Chatbots & Virtual Assistants — 24/7 student support, reduced admin workload",
    "     (e.g., Georgia State Univ. \"Pounce\" reduced enrollment melt by 21%)",
    "c) Automated Assessment & Feedback — NLP-based essay grading, auto-graded coding",
    "d) Predictive Analytics — attendance/grades/engagement data → targeted interventions",
])

make_two_column_slide(
    "AI — Impacts, Benefits & Costs",
    "Benefits",
    [
        "Student retention improvement: 10–15%",
        "Admin cost reduction: SGD 2–4M/year",
        "New CET revenue: SGD 5–10M/year",
        "24/7 chatbot support for students",
        "Faster assignment feedback",
        "AI-powered career guidance / job matching",
    ],
    "Cost-Benefit & ROI",
    [
        "Year 1 investment: ~SGD 4.5–6.5M",
        "Breakeven in Year 2",
        "5-year net return: ~SGD 15.6M (226% ROI)",
        "",
        "Key challenges:",
        "  • PDPA data privacy compliance",
        "  • Change management for lecturers",
        "  • Data quality & governance",
        "  • Integration with existing LMS/IT",
    ]
)

# 8 ── VR/AR ────────────────────────────────────────────────────────────────
make_section_slide("Technology 2: Immersive Media (VR/AR) — Collin")

make_content_slide("VR/AR — How It Improves NYP", [
    "a) Virtual Laboratories — safe, repeatable experiments (engineering, healthcare, science)",
    "     (e.g., Labster — 3,000+ institutions, 200+ virtual lab experiments)",
    "b) Immersive Learning — AR overlays of 3D models (anatomy, architecture, machines)",
    "     (e.g., Microsoft HoloLens in medical schools)",
    "c) Virtual Campus Tours — 360° tours for prospective & international students",
    "d) Industry Training Simulations — aviation, construction, hospitality training for CET",
    "     Licensable to industry partners → new revenue stream",
])

make_two_column_slide(
    "VR/AR — Impacts, Benefits & Costs",
    "Benefits",
    [
        "Expanded practical training capacity",
        "Equipment maintenance savings: 15–25%",
        "B2B licensing revenue: SGD 2–5M/year",
        "75% better knowledge retention vs. traditional (PwC)",
        "Enhanced student recruitment",
        "Differentiation from competitors",
    ],
    "Cost-Benefit & ROI",
    [
        "Year 1 investment: ~SGD 2.3–4.1M",
        "Breakeven in Year 2",
        "5-year net return: ~SGD 10.3M (198% ROI)",
        "",
        "Key challenges:",
        "  • High upfront hardware/content costs",
        "  • Specialized content development skills",
        "  • Motion sickness / accessibility",
        "  • Hardware refresh every 2–3 years",
        "  • Pedagogical integration training",
    ]
)

# 9 ── 5G & Edge Computing ─────────────────────────────────────────────────
make_section_slide("Technology 3: 5G & Edge Computing — Alex")

make_content_slide("5G & Edge Computing — How It Improves NYP", [
    "a) Smart Campus Infrastructure",
    "     IoT sensor networks, smart energy management (10–15% savings), smart parking",
    "     (e.g., NUS achieved 15% energy savings with IoT + edge computing)",
    "b) Enhanced Online & Hybrid Learning",
    "     4K/8K streaming, real-time collaboration, live remote lab access",
    "c) Enabler for VR/AR & AI",
    "     High bandwidth + low latency for VR; edge servers for real-time AI inference",
    "d) Campus Security & Safety",
    "     AI-powered surveillance, automated access control, emergency alerts",
])

make_two_column_slide(
    "5G & Edge — Impacts, Benefits & Costs",
    "Benefits",
    [
        "Modernized campus infrastructure",
        "Energy savings: 10–15%",
        "Foundation for VR/AR, AI, IoT",
        "Seamless online/hybrid learning",
        "5G speeds: up to 20 Gbps, ≤1 ms latency",
        "1M devices/km² connectivity",
    ],
    "Cost-Benefit & ROI",
    [
        "Year 1 investment: ~SGD 5–7M",
        "Breakeven in Year 3–4",
        "5-year net return: ~SGD 4.8M (67% ROI)",
        "(Infrastructure investment that accelerates",
        "  ROI of the other 3 technologies)",
        "",
        "Key challenges:",
        "  • High infrastructure CAPEX",
        "  • Spectrum allocation (IMDA)",
        "  • Expanded cybersecurity attack surface",
        "  • Specialized IT skills required",
    ]
)

# 10 ── Blockchain ──────────────────────────────────────────────────────────
make_section_slide("Technology 4: Blockchain Certificate Verification — Declan")

make_content_slide("Blockchain — How It Improves NYP", [
    "a) Digital Diploma & Certificate Issuance",
    "     Cryptographically signed, blockchain-stored, QR/URL-verifiable credentials",
    "     Builds on Singapore's OpenCerts ecosystem (opencerts.io)",
    "b) Micro-Credentialing & Skill Badges",
    "     Stackable, shareable on LinkedIn, granular skill portfolios",
    "     (e.g., MIT Blockcerts since 2017)",
    "c) Streamlined Verification — from 3–10 days → seconds",
    "     Eliminates manual admin processing",
    "d) Transcript & Academic Record Management",
    "     Selective sharing, credit transfer, SkillsFuture integration",
])

make_two_column_slide(
    "Blockchain — Impacts, Benefits & Costs",
    "Benefits",
    [
        "Manual verification workload: −50–70%",
        "Tamper-proof credentials (fraud eliminated)",
        "Verification services revenue: SGD 0.5–1M/year",
        "Instant employer verification (seconds)",
        "Portable, cross-border credentials",
        "Aligned with OpenCerts & Smart Nation",
    ],
    "Cost-Benefit & ROI",
    [
        "Year 1 investment: ~SGD 0.8–1.5M",
        "Breakeven in Year 3",
        "5-year net return: ~SGD 1.59M (88% ROI)",
        "(Low absolute cost; high strategic value)",
        "",
        "Key challenges:",
        "  • Employer / institution adoption",
        "  • Blockchain platform interoperability",
        "  • PDPA data privacy (off-chain PII)",
        "  • Scalability of blockchain network",
    ]
)

# 11 ── Consolidated Investment ─────────────────────────────────────────────
make_section_slide("Consolidated Investment & ROI")

make_content_slide("Consolidated Technology Investment Summary", [
    "Total Year 1 Investment: SGD 8.8–17.1 million",
    "Projected Annual Returns: SGD 10.5–23 million (revenue + cost savings)",
    "",
    "Individual Technology ROI:",
    "   • AI:         5-yr net ~SGD 15.6M  (226% ROI) — Breakeven Yr 2",
    "   • VR/AR:      5-yr net ~SGD 10.3M  (198% ROI) — Breakeven Yr 2",
    "   • 5G/Edge:    5-yr net ~SGD 4.8M   ( 67% ROI) — Breakeven Yr 3–4",
    "   • Blockchain:  5-yr net ~SGD 1.59M  ( 88% ROI) — Breakeven Yr 3",
    "",
    "Overall breakeven within 1–2 years for most initiatives.",
])

# 12 ── Roadmap ─────────────────────────────────────────────────────────────
make_section_slide("Implementation Roadmap")

make_content_slide("Three-Phase Implementation Roadmap", [
    "Phase 1 — Year 1: Foundation",
    "   • Deploy AI chatbot for student services",
    "   • Implement blockchain credential system (on OpenCerts)",
    "   • Begin 5G infrastructure planning & pilot",
    "   • Develop first VR lab prototypes",
    "",
    "Phase 2 — Year 2: Expansion",
    "   • Scale AI adaptive learning across all schools",
    "   • Launch 10+ VR virtual labs",
    "   • Deploy 5G campus-wide",
    "   • Expand blockchain to micro-credentials",
    "",
    "Phase 3 — Year 3: Optimization",
    "   • Integrate all technologies into cohesive smart campus",
    "   • Monetize VR content & blockchain verification services",
    "   • Comprehensive ROI evaluation & strategy adjustment",
])

# 13 ── Technology Synergy ──────────────────────────────────────────────────
make_section_slide("Technology Synergy Ecosystem")

make_content_slide("How the Four Technologies Reinforce Each Other", [
    "5G enables VR/AR — High-bandwidth, low-latency 5G is essential for seamless VR/AR.",
    "",
    "AI enhances VR/AR — AI personalizes VR/AR learning scenarios per student performance.",
    "",
    "Blockchain secures AI credentials — As AI-driven assessment grows, blockchain ensures",
    "   integrity and verifiability of AI-generated evaluations.",
    "",
    "Edge Computing supports all — Edge servers provide the low-latency processing needed",
    "   by AI inference, VR rendering, and IoT data processing.",
    "",
    "Together they form a cohesive smart-campus ecosystem.",
])

# 14 ── Conclusion ──────────────────────────────────────────────────────────
make_section_slide("Conclusion")

make_content_slide("Conclusion", [
    "AI, VR/AR, 5G/Edge Computing, and Blockchain can strategically transform NYP:",
    "",
    "   • Address current operational & educational challenges",
    "   • Enhance every stage of the educational value chain",
    "   • Generate SGD 10.5–23M/year in revenue + cost savings",
    "   • Position NYP as Singapore's leading digital-forward polytechnic",
    "",
    "A phased 3-year roadmap manages risk while maximising synergies.",
    "",
    "The four technologies are not independent — they create a synergistic ecosystem",
    "that reinforces each solution and delivers compounding benefits.",
])

# 15 ── Thank You ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
            "Thank You", font_size=44, bold=True, colour=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(2),
            "Questions & Discussion\n\n"
            "Team: Clifton (AI) • Collin (VR/AR) • Alex (5G & Edge) • Declan (Blockchain)",
            font_size=20, colour=RGBColor(0xBB, 0xCC, 0xDD),
            alignment=PP_ALIGN.CENTER)

# ── save ────────────────────────────────────────────────────────────────────
output_path = "presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
