# Requirements
# Install with: pip install -r requirements.txt

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.core_properties.title = "Haidilao Smart Table — AI Health Assistant"
prs.core_properties.author = "Team 1 Clifton"

slides_content = [
    {"title": "Haidilao “Smart Table” — AI Health Assistant for Dining", "bullets": ["Real-time food intake tracking via wristbands + camera"], "notes": "Title slide. One-line hook: ‘Premium dining meets personalised health — live at the table.’"},
    {"title": "Introduction & Company background", "bullets": ["Haidilao: premium hotpot known for exceptional service and tech-forward experiences.", "New service: AI-powered food intake tracking — camera + wristbands + app.", "Core benefit: personalised, real-time dietary feedback that enhances the dining experience and supports health goals."], "notes": "Frame the service as a service-layer extension of Haidilao’s ‘service-first’ brand."},
    {"title": "Market opportunity & trends", "bullets": ["Rising health & wellness dining trend among urban millennials and Gen Z.", "Growth in personalised nutrition and appetite for tech-enabled experiences.", "Consumers willing to pay premium for value-added, privacy-conscious services."], "notes": "Cite trend: health-tech adoption + experiential dining = clear product-market fit."},
    {"title": "How it works — user flow", "bullets": ["1. Scan QR on table via Haidilao app.", "2. Choose daily health goal + register wristband color.", "3. Overhead camera maps wristbands to food items taken.", "4. Real-time tips + end-of-meal nutrition summary in meal history."], "notes": "Keep flow simple for demos; emphasise privacy (no facial ID; wristband ID only)."},
    {"title": "Fit with Haidilao brand & USP", "bullets": ["Builds on ‘service excellence’: proactive, personalised guidance at the moment of choice.", "USP: only major hotpot chain delivering real-time, table-level nutrition guidance integrated with loyalty/profile history.", "Differentiator: combination of hospitality + AI health assistant (not calorie-shaming — supportive)."], "notes": "Emphasise emotional value — guests feel cared for and entertained."},
    {"title": "Competitor 1 — Little Sheep (Xiao Fei Yang) — SWOT", "bullets": ["Strengths: Established hotpot brand; widespread locations; price-competitive.", "Weaknesses: Less premium, fewer tech-driven experiences; lower loyalty integration.", "Opportunities: Could adopt health messaging with simpler menu cues.", "Threats: Haidilao’s superior service & tech could pull health-conscious premium customers."], "notes": "Little Sheep is strong on footprint but weak on experiential tech — opportunity for Haidilao to own premium health-tech."},
    {"title": "Competitor 2 — Xiabu Xiabu — SWOT", "bullets": ["Strengths: Fast service, low price; appeals to value diners.", "Weaknesses: Commodity perception; limited app/loyalty features.", "Opportunities: Partnership with delivery/meal-tracking apps for lower-cost health features.", "Threats: New entrants offering healthy-fast hotpot options; regulatory privacy concerns could affect camera/AI adoption."], "notes": "Competitors may follow but need time and capital to replicate integrated in-dining experience."},
    {"title": "Marketing objectives (SMART)", "bullets": ["Adoption: 20% adoption among Haidilao app users dining at pilot stores within 6 months.", "Engagement: 40% repeat usage within 90 days.", "Monetary: Increase average per-person ancillary revenue by 6% in pilot locations.", "Brand: Generate 5,000 social impressions from organic + influencer posts in first 3 months."], "notes": "KPIs: adoption rate, DAU in-app, repeat rate, ARPU/ANC revenue, impressions."},
    {"title": "Target market & segmentation", "bullets": ["Primary: Urban, health-interested diners 22–45 (millennials + Gen Z), dual-income, tech-savvy.", "Segmentation: Demographic, Psychographic (health-conscious, tech adopters), Behavioral (monthly 2+ visits, loyalty members).", "Pilot audience: loyalty members and weekday dinner groups in Tier-1 cities."], "notes": "Focus pilots on high-LTV customers who value both experience and health insights."},
    {"title": "Product strategy", "bullets": ["Core: wristband + overhead camera detection + in-app real-time feedback + meal history.", "Tiers: Basic (free summary), Plus (detailed nutrition + weekly summary), Premium (nutrition coaching + diet integration).", "Core benefit: immediate, personalised dietary guidance that enhances enjoyment without judgement."], "notes": "Emphasise low-friction onboarding, privacy-first design (camera recognizes wristband color, not faces)."},
    {"title": "Pricing strategy", "bullets": ["Freemium launch: Basic free during pilot.", "Intro paid tier: Plus at small fee or single-meal add-on after trial.", "Loyalty integration: redeem points for Plus/Premium; VIP free trial.", "Tactics: limited-time discounts, bundle with dessert/drink vouchers, A/B test price points."], "notes": "Keep initial friction low: freemium -> trial -> paid. Use loyalty mechanics to reduce acquisition cost."},
    {"title": "Promotion strategy", "bullets": ["In-restaurant: table tents, staff demos, QR stickers, wristband display.", "App: push notifications, onboarding flow, in-app banners.", "Social & PR: short demo reels, influencer dining experiences, health KOLs.", "Partnerships: health apps, corporate wellness programs."], "notes": "Emphasise experiential promotions and shareability — wristbands are fun social props."},
    {"title": "Distribution (Place) strategy", "bullets": ["Pilot: 10–15 flagship stores in Tier-1 cities.", "Scale: roll out to premium outlets in 9–12 months based on pilot metrics.", "Channel: App-first experience; on-premise hardware; franchise training & SOPs."], "notes": "Use owned channels (restaurants + app) for full control; franchise standards ensure consistent experience."},
    {"title": "Conclusion & Further recommendations", "bullets": ["Aligns with Haidilao’s service promise; adds personal care and entertainment.", "Extensions: privacy-first certifications, integrate with wearables, gamify table challenges, corporate wellness packages, modular hardware for pop-ups."], "notes": "Reinforce measurable objectives and service’s ability to deepen customer relationships."}
]

for slide_data in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = slide_data['title']
    body.clear()
    for i, bullet in enumerate(slide_data['bullets']):
        if i == 0:
            p = body.paragraphs[0]
            p.text = bullet
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = slide_data.get('notes', '')

output_filename = "Haidilao_SmartTable_MarketingPlan.pptx"
prs.save(output_filename)
print(f"Saved presentation to {output_filename}")
